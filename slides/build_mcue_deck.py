from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
summary_path = ROOT / "results" / "sampled_summary.json"
if not summary_path.exists():
    summary_path = ROOT / "outputs" / "mcue_modern" / "summary.json"
SUMMARY = json.loads(summary_path.read_text(encoding="utf-8"))
OUT = ROOT / "slides" / "output" / "mcue_group_meeting_deck.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

COLORS = {
    "bg": RGBColor(0xFF, 0xF8, 0xEE),
    "panel": RGBColor(0xF5, 0xEB, 0xDE),
    "ink": RGBColor(0x10, 0x20, 0x33),
    "muted": RGBColor(0x5A, 0x68, 0x7A),
    "accent": RGBColor(0xB8, 0x4E, 0x2E),
    "accent2": RGBColor(0xE3, 0xA1, 0x3B),
    "line": RGBColor(0xC6, 0xB6, 0xA1),
}


def add_bg(slide, color=COLORS["bg"]):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(10.8), Inches(-0.2), Inches(4.3), Inches(8.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF6, 0xE3, 0xC7)
    shape.fill.transparency = 0.2
    shape.line.color.rgb = RGBColor(0xF6, 0xE3, 0xC7)
    shape.rotation = -8


def textbox(slide, left, top, width, height, text, size=20, color=COLORS["ink"], bold=False, font="Aptos", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def bullet_box(slide, left, top, width, height, lines, size=18):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = COLORS["ink"]
        p.bullet = True
    return box


def add_picture_frame(slide, image_path: Path, left, top, width, height, label: str | None = None):
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    frame.line.color.rgb = COLORS["line"]
    frame.line.width = Pt(1.2)
    slide.shapes.add_picture(str(image_path), left + Inches(0.08), top + Inches(0.08), width - Inches(0.16), height - Inches(0.26))
    if label:
        textbox(slide, left + Inches(0.08), top + height - Inches(0.18), width - Inches(0.16), Inches(0.16), label, size=11, color=COLORS["muted"], align=PP_ALIGN.CENTER)


def metric(v: float) -> str:
    return f"{v:.3f}"


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
textbox(slide, Inches(0.75), Inches(0.65), Inches(5.9), Inches(0.8), "Modernizing Magnetic Tile Saliency Detection", size=27, bold=True, font="Georgia")
textbox(
    slide,
    Inches(0.78),
    Inches(1.55),
    Inches(6.2),
    Inches(0.6),
    "From the legacy VS2013 toolbox to an auto-runnable Python MCue pipeline",
    size=18,
    color=COLORS["muted"],
)
line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.8), Inches(2.45), Inches(1.5), Inches(0.03))
line.fill.solid()
line.fill.fore_color.rgb = COLORS["accent2"]
line.line.color.rgb = COLORS["accent2"]
textbox(
    slide,
    Inches(0.82),
    Inches(2.8),
    Inches(5.2),
    Inches(1.1),
    "Dataset: author's Magnetic-tile-defect-datasets\nBatch run: 24 sampled images (4 per class)",
    size=18,
    color=COLORS["ink"],
)
textbox(
    slide,
    Inches(0.82),
    Inches(4.35),
    Inches(5.5),
    Inches(0.8),
    f"Overall sampled metrics\nF1 {metric(SUMMARY['overall']['f1'])}   IoU {metric(SUMMARY['overall']['iou'])}   MAE {metric(SUMMARY['overall']['mae'])}",
    size=18,
    bold=True,
    color=COLORS["accent"],
)
textbox(slide, Inches(0.82), Inches(6.7), Inches(2.5), Inches(0.3), "Zoom group meeting version", size=13, color=COLORS["muted"])
slide.shapes.add_picture(str(ROOT / "git-magnetic-tile-datasets" / "dataset.png"), Inches(6.15), Inches(1.0), Inches(6.3), Inches(4.9))


# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLORS["bg"])
textbox(slide, Inches(0.75), Inches(0.55), Inches(5.8), Inches(0.6), "What Was Ported From The Original Code", size=24, bold=True, font="Georgia")
textbox(
    slide,
    Inches(0.78),
    Inches(1.35),
    Inches(5.8),
    Inches(0.8),
    "Core MCue fusion preserved from McueSalTest.cpp:\nDarker prior + structural tensor + PHOT + AC + BMS",
    size=17,
    color=COLORS["ink"],
)
bullet_box(
    slide,
    Inches(0.85),
    Inches(2.4),
    Inches(5.8),
    Inches(2.8),
    [
        "Python + OpenCV headless instead of VS2013/OpenCV 3.1",
        "Direct use of paired JPG image and PNG mask files from the dataset repo",
        "Automatic export of per-cue maps, binary masks, metrics, and presentation-ready panels",
    ],
    size=17,
)
cue_dir = ROOT / "outputs" / "mcue_modern" / "MT_Blowhole" / "exp1_num_108889"
cue_specs = [
    ("darker.png", "Darker prior"),
    ("tensor.png", "Tensor"),
    ("phot.png", "PHOT"),
    ("ac.png", "AC"),
    ("bms.png", "BMS"),
    ("mcue.png", "Final MCue"),
]
for idx, (file_name, label) in enumerate(cue_specs):
    col = idx % 2
    row = idx // 2
    left = Inches(7.0 + col * 2.8)
    top = Inches(1.35 + row * 1.9)
    add_picture_frame(slide, cue_dir / file_name, left, top, Inches(2.45), Inches(1.55), label)


# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0xFA, 0xF6, 0xEF))
textbox(slide, Inches(0.75), Inches(0.55), Inches(4.8), Inches(0.6), "Sampled Results On 24 Images", size=24, bold=True, font="Georgia")
textbox(slide, Inches(0.8), Inches(1.28), Inches(6.8), Inches(0.35), "4 images per class, binary mask derived by Otsu thresholding on the saliency map", size=15, color=COLORS["muted"])

headers = ["Class", "F1", "IoU", "MAE"]
classes = ["MT_Blowhole", "MT_Break", "MT_Crack", "MT_Fray", "MT_Free", "MT_Uneven"]
x_positions = [0.82, 2.75, 3.6, 4.45]
for header, x in zip(headers, x_positions):
    textbox(slide, Inches(x), Inches(1.95), Inches(0.9 if header != "Class" else 1.6), Inches(0.25), header, size=14, color=COLORS["accent"], bold=True)
rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.82), Inches(2.28), Inches(4.35), Inches(0.02))
rule.fill.solid()
rule.fill.fore_color.rgb = COLORS["line"]
rule.line.color.rgb = COLORS["line"]
for idx, cls in enumerate(classes):
    y = 2.45 + idx * 0.45
    metrics = SUMMARY["categories"][cls]
    textbox(slide, Inches(0.82), Inches(y), Inches(1.6), Inches(0.2), cls.replace("MT_", ""), size=14)
    textbox(slide, Inches(2.75), Inches(y), Inches(0.7), Inches(0.2), metric(metrics["f1"]), size=14, color=COLORS["ink"])
    textbox(slide, Inches(3.6), Inches(y), Inches(0.7), Inches(0.2), metric(metrics["iou"]), size=14, color=COLORS["ink"])
    textbox(slide, Inches(4.45), Inches(y), Inches(0.7), Inches(0.2), metric(metrics["mae"]), size=14, color=COLORS["ink"])

textbox(slide, Inches(6.1), Inches(1.85), Inches(2.8), Inches(0.35), "Main takeaways", size=21, bold=True, font="Georgia")
bullet_box(
    slide,
    Inches(6.15),
    Inches(2.45),
    Inches(5.8),
    Inches(3.2),
    [
        "Blowhole and crack are the strongest qualitative categories in the sampled run.",
        "Uneven and fray remain difficult because saliency spreads across textured regions.",
        "Free samples stay near-zero after thresholding, which is a useful sanity check.",
        "The port is faithful to the feature logic, but not a bitwise rebuild of the old executable.",
    ],
    size=17,
)
box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.15), Inches(5.95), Inches(4.7), Inches(0.95))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
box.line.color.rgb = COLORS["line"]
textbox(
    slide,
    Inches(6.38),
    Inches(6.18),
    Inches(4.2),
    Inches(0.5),
    f"Overall   F1 {metric(SUMMARY['overall']['f1'])}  |  IoU {metric(SUMMARY['overall']['iou'])}  |  MAE {metric(SUMMARY['overall']['mae'])}",
    size=17,
    color=COLORS["accent"],
    bold=True,
)


# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0xFF, 0xF7, 0xEF))
textbox(slide, Inches(0.75), Inches(0.55), Inches(4.5), Inches(0.6), "Qualitative Examples", size=24, bold=True, font="Georgia")
textbox(
    slide,
    Inches(0.8),
    Inches(1.25),
    Inches(7.5),
    Inches(0.35),
    "Each panel shows: original image | ground truth | predicted saliency | thresholded overlay",
    size=15,
    color=COLORS["muted"],
)
examples = [
    ("Blowhole | best sampled example", ROOT / "outputs" / "mcue_modern" / "MT_Blowhole" / "exp1_num_108889" / "panel.png"),
    ("Crack | stable localization", ROOT / "outputs" / "mcue_modern" / "MT_Crack" / "exp1_num_249594" / "panel.png"),
    ("Uneven | hard diffuse-texture case", ROOT / "outputs" / "mcue_modern" / "MT_Uneven" / "exp0_num_461" / "panel.png"),
]
for idx, (label, img) in enumerate(examples):
    top = 1.9 + idx * 1.75
    textbox(slide, Inches(0.84), Inches(top), Inches(4.0), Inches(0.22), label, size=14, color=COLORS["accent"], bold=True)
    add_picture_frame(slide, img, Inches(0.84), Inches(top + 0.28), Inches(11.4), Inches(1.25))

textbox(
    slide,
    Inches(0.85),
    Inches(6.95),
    Inches(11.6),
    Inches(0.3),
    "Next step for a stronger benchmark: sweep thresholds or use continuous saliency metrics, then scale the full run on your school SSH machine.",
    size=14,
    color=COLORS["ink"],
)

prs.save(str(OUT))
print(OUT)
